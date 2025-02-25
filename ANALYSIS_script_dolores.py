import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from scipy.optimize import curve_fit
from scipy.signal import find_peaks
from scipy.fftpack import rfft, irfft, fftfreq, fft
from mpl_toolkits.axes_grid1.inset_locator import inset_axes
from matplotlib.gridspec import GridSpec
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import math
import json

def add_image_with_aspect_ratio(slide, image_path, left, top, cell_width, cell_height):
    img = Image.open(image_path)
    img_width, img_height = img.size
    
    # Calculate aspect ratio
    aspect_ratio = img_width / img_height
    
    # Calculate new dimensions preserving aspect ratio
    if aspect_ratio > 1:  # Landscape
        new_width = cell_width
        new_height = cell_width / aspect_ratio
    else:  # Portrait
        new_height = cell_height
        new_width = cell_height * aspect_ratio
    
    # Center the image within the cell
    centered_left = left + (cell_width - new_width) / 2
    centered_top = top + (cell_height - new_height) / 2
    
    # Add the picture to the slide
    slide.shapes.add_picture(image_path, Inches(centered_left), Inches(centered_top),
                           width=Inches(new_width), height=Inches(new_height))

def ten_image_slides(prs, image_paths):
    """
    Add images to presentation with 10 images per slide until all images are added.
    """
    for k in range(math.ceil(len(image_paths)/15)):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Slide dimensions (inches)
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        # Calculate grid
        rows = 3
        cols = 5
        cell_width = slide_width / cols
        cell_height = slide_height / rows
        
        # Add images to slide
        for i in range(rows):
            for j in range(cols):
                idx = i * cols + j + k*15
                if idx < len(image_paths):
                    left = j * cell_width
                    top = i * cell_height
                    add_image_with_aspect_ratio(slide, image_paths[idx], left, top,
                                              cell_width, cell_height)

def most_common(lst):
    """Return most common element in a list."""
    return max(set(lst), key=lst.count)

def lorentzian(x, A, x0, gamma, B):
    """Calculate the Lorentzian function."""
    return ((A /np.pi) * (gamma / ((x - x0)**2 + gamma**2)) + B)

def analyze_data(DATA_DIRECTORY, TITLE, AUTHOR, wafer, reticle, chip, PD_Channel=1):
    chip_ID = f"{wafer}_{reticle}_{chip}"
    image_paths = []
    
    # Create columns of final data table
    f_Q_values = []
    f_resonances = []
    f_exctinction_ratios = []
    f_dQ_values = []
    f_params1 = []
    f_covar1 = []
    files = []
    
    for filename in os.listdir(DATA_DIRECTORY):
        if not filename.endswith('.csv'):
            continue
            
        T_values = []
        fig = plt.figure(figsize=(8, 10), layout="constrained")
        gs = GridSpec(3, 2, figure=fig)
        ax1 = fig.add_subplot(gs[0, :])
        ax3 = fig.add_subplot(gs[1, :])
        
        # Load and process data
        path = os.path.join(DATA_DIRECTORY, filename)
        df = pd.read_csv(path)
        x = np.array(df['wavelength'])
        y = np.array(df[f'channel_{PD_Channel}'])
        y_linear = y
        
        # Convert wavelength to frequency
        speed_of_light = 299792458
        x_frequency = [(1/x_val) * speed_of_light for x_val in x]
        
        # Find peaks
        height = np.mean(y_linear) + 1.5 * np.std(y_linear)
        peaks, properties = find_peaks(y_linear, height=height, distance=300)
        
        # Process peaks and calculate period
        x_coordinates = [x_frequency[peak] for peak in peaks]
        if len(x_coordinates) > 1:
            for i in range(len(x_coordinates)-1):
                frequency = abs(x_coordinates[i+1] - x_coordinates[i])
                T_values.append(1/frequency)
            T_values = [round(T_value, 3) for T_value in T_values]
            T = most_common(T_values)
            
            # FFT Processing
            N = len(y_linear)
            yf = fft(y_linear)
            xf = np.linspace(0, 1.0/(2.0*T), N//2)
            f_signal = rfft(y_linear)
            W = fftfreq(y_linear.size, d=T)
            
            # Filter signal
            cut_f_signal = f_signal.copy()
            cut_f_signal[W<0.2] = 0
            cut_f_signal[W>0.3] = 0
            cut_signal = irfft(cut_f_signal)
            
            # Process filtered data
            new_values = np.array(y_linear - cut_signal).flatten()
            offset = np.min(y_linear) - np.min(new_values)
            final_data = [filter_value + (offset if offset > 0 else -offset)
                         for filter_value in new_values]
            
            # Find peaks in filtered data
            peaks, properties = find_peaks(-np.array(final_data),
                                        height=2*np.std(-np.array(final_data))+
                                        np.mean(-np.array(final_data)),
                                        distance=500)
            
            # Process each peak
            Q_values = []
            x_peak_location = []
            fit_parameters = []
            dQ_values = []
            covar1 = []
            resonances = []
            exctinction_ratios = []
            
            for peak in peaks:
                num = 100
                xloc = x[peak-num:peak+num]
                yloc = final_data[peak-num:peak+num]
                
                try:
                    # Fit Lorentzian
                    p0 = [-(np.max(yloc)-np.min(yloc)), xloc[round(len(xloc)/2)],
                          (xloc[-1]-xloc[0])/2, 0.0001]
                    params, covar = curve_fit(lorentzian, xloc, yloc, p0=p0,
                                            bounds=((-np.inf, -np.inf, -np.inf, -np.inf),
                                                   (0, np.inf, np.inf, np.inf)))
                    
                    # Calculate metrics
                    Q = abs(params[1])/(abs(params[2])*2)
                    perr = np.sqrt(np.diag(covar))
                    dQ = 100*((np.sqrt(((perr[1]/params[1])**2)+((perr[2]/params[2])**2))))
                    resonance = params[1]
                    extinction_ratio = 10*np.log10(min(yloc)/max(yloc))
                    
                    # Filter results
                    if (x[peak] < 1615 and x[peak] > 1510 and dQ < 10 and
                        Q > 4000 and extinction_ratio < (-1)):
                        Q_values.append(Q)
                        x_peak_location.append(xloc)
                        fit_parameters.append(params)
                        dQ_values.append(dQ)
                        covar1.append(covar)
                        resonances.append(resonance)
                        exctinction_ratios.append(extinction_ratio)
                
                except (RuntimeError, ValueError):
                    continue
            
            # Plot and save results if peaks found
            if len(Q_values) > 0:
                # Plot processing and save figure
                Q_max = np.max(Q_values)
                index_Q = Q_values.index(Q_max)
                x_index = np.where(x == x_peak_location[index_Q][0])[0][0]
                
                # Save plot if Q > 10000
                if Q_max > 10000:
                    image_paths.append(os.path.join(DATA_DIRECTORY,
                                                  filename.replace('.csv', '.png')))
                
                # Save results
                f_Q_values.append(Q_values)
                f_resonances.append(resonances)
                f_exctinction_ratios.append(exctinction_ratios)
                f_dQ_values.append(dQ_values)
                f_params1.append(fit_parameters)
                f_covar1.append(covar1)
                files.append(filename)
    
    # Create results dictionary
    results = {
        chip_ID: {
            'device name': files,
            'Quality Factor': f_Q_values,
            'Error in Quality Factor (%)': f_dQ_values,
            'Resonance (nm)': f_resonances,
            'Exctinction ratio (dB)': f_exctinction_ratios
        }
    }
    
    return results, image_paths

def create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH, SUMMARY_POINTS=''):
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TITLE
    subtitle.text = AUTHOR
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = 'Summary'
    content.text = SUMMARY_POINTS
    
    # Add image slides
    ten_image_slides(prs, image_paths)
    
    # Save presentation
    prs.save(PPTX_SAVE_PATH)

def main():
    # Configuration
    TITLE = 'SOF2_C4_11'
    AUTHOR = 'Ana Elias & Skylar Deckoff-Jones'
    wafer = "SOF2"
    reticle = "D3"
    chip = "52"
    DATA_DIRECTORY = 'path/to/your/data/directory'
    PPTX_SAVE_PATH = 'path/to/save/presentation.pptx'
    JSON_SAVE_PATH = 'path/to/save/results.json'
    PD_Channel = 1
    
    # Run analysis
    results, image_paths = analyze_data(DATA_DIRECTORY, TITLE, AUTHOR,
                                      wafer, reticle, chip, PD_Channel)
    
    # Save results to JSON
    with open(JSON_SAVE_PATH, 'w') as json_file:
        json.dump(results, json_file, indent=4)
    
    # Create presentation
    create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH)

if __name__ == "__main__":
    main()
