import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from scipy.optimize import curve_fit
from scipy.signal import find_peaks
from scipy.fftpack import rfft, irfft, fftfreq, fft
from mpl_toolkits.axes_grid1.inset_locator import inset_axes
from matplotlib.gridspec import GridSpec
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import math
import json
import re

def extract_voltage_from_filename(filename):
    """Extract voltage value from filename."""
    match = re.search(r'Voltage_(-?\d+\.\d+)V_', filename)
    if match:
        return float(match.group(1))
    return None

def most_common(lst):
    """Return most common element in a list."""
    return max(set(lst), key=lst.count)

def lorentzian(x, A, x0, gamma, B):
    """Calculate the Lorentzian function."""
    return ((A /np.pi) * (gamma / ((x - x0)**2 + gamma**2)) + B)

def plot_transmission_spectra(data_directory, channel=1, save_dir=None):
    """Create overlay plot of transmission spectra at different voltages."""
    print("\nCreating transmission spectra overlay plot...")
    
    plt.figure(figsize=(12, 6))
    
    # Get all CSV files and sort by voltage
    csv_files = [(f, extract_voltage_from_filename(f))
                 for f in os.listdir(data_directory)
                 if f.endswith('.csv')]
    csv_files = [(f, v) for f, v in csv_files if v is not None]
    csv_files.sort(key=lambda x: x[1])
    
    # Create color map for different voltages
    colors = plt.cm.rainbow(np.linspace(0, 1, len(csv_files)))
    
    for (filename, voltage), color in zip(csv_files, colors):
        # Load data
        df = pd.read_csv(os.path.join(data_directory, filename))
        wavelength = df['wavelength']
        transmission = df[f'channel_{channel}']
        
        # Convert to dBm if not already
        if transmission.max() > 0:  # assuming linear scale
            transmission = 10 * np.log10(transmission)
            
        plt.plot(wavelength, transmission, label=f'{voltage}V', color=color, alpha=0.7)
    
    plt.xlabel('Wavelength (nm)')
    plt.xlim(1526.2, 1526.35)  # or whatever range your resonances are in
    plt.ylabel('Transmission (dBm)')
    plt.title('Transmission Spectra at Different Voltages')
    plt.grid(True)
    plt.legend(bbox_to_anchor=(1.05, 1), loc='upper left')
    plt.tight_layout()
    
    # Save plot
    plot_path = os.path.join(save_dir, 'transmission_spectra.png')
    plt.savefig(plot_path, bbox_inches='tight')
    plt.close()
    
    print(f"Transmission spectra plot saved to: {plot_path}")
    return plot_path

def analyze_voltage_series(data_directory, channel=1):
    """Analyze voltage-dependent measurements from CSV files."""
    print(f"\nStarting analysis for channel_{channel}")
    print(f"Looking in directory: {data_directory}")
    
    voltage_data = {
        'voltages': [],
        'resonance_wavelengths': [],
        'q_factors': [],
        'dbm_values': [],
        'peak_indices': []
    }
    
    # Get all CSV files and sort by voltage
    csv_files = [(f, extract_voltage_from_filename(f))
                 for f in os.listdir(data_directory)
                 if f.endswith('.csv')]
    csv_files = [(f, v) for f, v in csv_files if v is not None]
    csv_files.sort(key=lambda x: x[1])
    
    print(f"\nFound {len(csv_files)} CSV files")
    for filename, voltage in csv_files:
        print(f"\nProcessing: {filename} (Voltage: {voltage}V)")
        
        try:
            # Load and process data
            path = os.path.join(data_directory, filename)
            df = pd.read_csv(path)
            print(f"DataFrame shape: {df.shape}")
            print(f"Columns: {df.columns.tolist()}")
            
            x = np.array(df['wavelength'])
            y = np.array(df[f'channel_{channel}'])
            print(f"Wavelength range: {x.min():.2f} to {x.max():.2f}")
            print(f"Signal range: {y.min():.2f} to {y.max():.2f}")
            
            # Initial peak finding
            height = np.mean(y) - 1.5 * np.std(y)  # Changed + to - to look for dips
            peaks, _ = find_peaks(-y, height=-height)
            print(f"Found {len(peaks)} initial peaks")
            
            if len(peaks) > 1:
                # Calculate period using frequency analysis
                speed_of_light = 299792458
                x_frequency = [(1/x_val) * speed_of_light for x_val in x]
                x_coordinates = [x_frequency[peak] for peak in peaks]
                
                T_values = []
                for i in range(len(x_coordinates)-1):
                    frequency = abs(x_coordinates[i+1] - x_coordinates[i])
                    T_values.append(1/frequency)
                
                if T_values:
                    T_values = [round(T_value, 3) for T_value in T_values]
                    T = most_common(T_values)
                    print(f"Calculated period: {T}")
                    
                    # FFT Processing
                    N = len(y)
                    f_signal = rfft(y)
                    W = fftfreq(y.size, d=T)
                    
                    # Filter signal
                    cut_f_signal = f_signal.copy()
                    cut_f_signal[W<0.2] = 0
                    cut_f_signal[W>0.3] = 0
                    cut_signal = irfft(cut_f_signal)
                    
                    # Process filtered data
                    new_values = np.array(y - cut_signal).flatten()
                    offset = np.min(y) - np.min(new_values)
                    final_data = [val + (offset if offset > 0 else -offset)
                                for val in new_values]
                    
                    # Find peaks in filtered data
                    peaks, _ = find_peaks(-np.array(final_data),
                                       height=2*np.std(-np.array(final_data))+np.mean(-np.array(final_data)),
                                       distance=500)
                    print(f"Found {len(peaks)} peaks after filtering")
                    
                    # Process each peak
                    for peak in peaks:
                        try:
                            num = 100
                            if peak-num < 0 or peak+num >= len(x):
                                continue
                                
                            xloc = x[peak-num:peak+num]
                            yloc = final_data[peak-num:peak+num]
                            
                            # Fit Lorentzian
                            p0 = [-(np.max(yloc)-np.min(yloc)),
                                 xloc[round(len(xloc)/2)],
                                 (xloc[-1]-xloc[0])/2,
                                 0.0001]
                            
                            params, covar = curve_fit(lorentzian, xloc, yloc, p0=p0,
                                                    bounds=((-np.inf, -np.inf, -np.inf, -np.inf),
                                                           (0, np.inf, np.inf, np.inf)))
                            
                            # Calculate metrics
                              # Calculate metrics
                            Q = abs(params[1])/(abs(params[2])*2)
                            resonance = params[1]

                            # Calculate extinction ratio (data already in dB)
                            raw_transmission = y[peak-num:peak+num] ##Lets come back to this please.
                            extinction_ratio = -(min(raw_transmission) - max(raw_transmission))  # Get positive value
                            
                            # Store results (with relaxed filtering)
                              # Store results with extinction ratio threshold
                            if 1526.2 < x[peak] < 1526.35 and abs(extinction_ratio) > 2.5:  # Only accept peaks with extinction ratio > 0.5 dB
                                voltage_data['voltages'].append(voltage)
                                voltage_data['resonance_wavelengths'].append(resonance)
                                voltage_data['q_factors'].append(Q)
                                voltage_data['dbm_values'].append(extinction_ratio)
                                voltage_data['peak_indices'].append(peak)
                                print("  Peak accepted")
                            else:
                                print("  Peak rejected (wavelength out of range or low extinction ratio)")
                                
                        except Exception as e:
                            print(f"Error processing peak: {str(e)}")
                            continue
                            
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            continue
    
    print("\nAnalysis complete:")
    print(f"Total peaks found: {len(voltage_data['voltages'])}")
    return voltage_data

def plot_voltage_dependencies(voltage_data, save_dir):
    """Create plots of parameters vs voltage."""
    print("\nCreating voltage dependency plots...")
    
    fig = plt.figure(figsize=(15, 5))
    
    # Resonance Wavelength vs Voltage
    ax1 = fig.add_subplot(131)
    ax1.scatter(voltage_data['voltages'], voltage_data['resonance_wavelengths'])
    ax1.set_xlabel('Voltage (V)')
    ax1.set_ylabel('Resonance Wavelength (nm)')
    ax1.set_title('Resonance Wavelength vs Voltage')
    ax1.grid(True)
    
    # dBm vs Voltage
    ax2 = fig.add_subplot(132)
    ax2.scatter(voltage_data['voltages'], voltage_data['dbm_values'])
    ax2.set_xlabel('Voltage (V)')
    ax2.set_ylabel('Extinction Ratio (dB)')
    ax2.set_title('Extinction Ratio vs Voltage')
    ax2.grid(True)
    
    # Q factor vs Voltage
    ax3 = fig.add_subplot(133)
    ax3.scatter(voltage_data['voltages'], voltage_data['q_factors'])
    ax3.set_xlabel('Voltage (V)')
    ax3.set_ylabel('Q Factor')
    ax3.set_title('Q Factor vs Voltage')
    ax3.grid(True)
    
    plt.tight_layout()
    
    plot_path = os.path.join(save_dir, 'voltage_dependencies.png')
    plt.savefig(plot_path)
    plt.close()
    
    print(f"Plots saved to: {plot_path}")
    return plot_path

def create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH, SUMMARY_POINTS=''):
    """Create PowerPoint presentation with results."""
    print("\nCreating PowerPoint presentation...")
    
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TITLE
    subtitle.text = AUTHOR
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = 'Voltage-Dependent Analysis'
    content.text = SUMMARY_POINTS or "Analysis of resonance parameters vs applied voltage"
    
    # Add plots
    if image_paths:
        slide_layout = prs.slide_layouts[5]
        
        # First plot - voltage dependencies
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(4)
        slide.shapes.add_picture(image_paths[0], left, top, width, height)
        
        # Second plot - transmission spectra
        if len(image_paths) > 1:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(image_paths[1], left, top, width, height)
    
    prs.save(PPTX_SAVE_PATH)
    print(f"Presentation saved to: {PPTX_SAVE_PATH}")

def analyze_data(DATA_DIRECTORY, TITLE, AUTHOR, wafer, reticle, chip, PD_Channel=1):
    """Main analysis function."""
    print(f"\nStarting analysis for {wafer}_{reticle}_{chip}")
    chip_ID = f"{wafer}_{reticle}_{chip}"
    
    # Analyze voltage series
    voltage_data = analyze_voltage_series(DATA_DIRECTORY, PD_Channel)
    
    # Create voltage dependency plots
    voltage_plot_path = plot_voltage_dependencies(voltage_data, DATA_DIRECTORY)
    
    # Create transmission spectra overlay plot
    spectra_plot_path = plot_transmission_spectra(DATA_DIRECTORY, PD_Channel, DATA_DIRECTORY)
    
    # Create results dictionary
    results = {
        chip_ID: {
            'voltages': voltage_data['voltages'],
            'resonance_wavelengths': voltage_data['resonance_wavelengths'],
            'q_factors': voltage_data['q_factors'],
            'extinction_ratios': voltage_data['dbm_values']
        }
    }
    
    return results, [voltage_plot_path, spectra_plot_path]

def main():
    """Main execution function."""
    # === CONFIGURATION - EDIT THESE VALUES ===
    TITLE = 'Voltage-Dependent Analysis'
    AUTHOR = 'aadarsh'
    wafer = "Sample"
    reticle = "Test"
    chip = "1"
    
    # Set this to the full path of your Sweep_Data folder
    DATA_DIRECTORY = '/Users/aadarsh/Desktop/memq/10V'
    
    # Where to save results
    PPTX_SAVE_PATH = 'voltage_analysis_results.pptx'
    JSON_SAVE_PATH = 'voltage_analysis_results.json'
    
    # Analysis settings
    PD_Channel = 1  # Which channel to analyze (1-4)
    # =====================================
    
    print("\nStarting voltage analysis script...")
    print(f"Data directory: {DATA_DIRECTORY}")
    print(f"Analyzing channel: {PD_Channel}")
    
    try:
        # Run analysis
        results, image_paths = analyze_data(DATA_DIRECTORY, TITLE, AUTHOR,
                                          wafer, reticle, chip, PD_Channel)
        
        # Save results to JSON
        with open(JSON_SAVE_PATH, 'w') as json_file:
            json.dump(results, json_file, indent=4)
        print(f"\nResults saved to: {JSON_SAVE_PATH}")
        
        # Create presentation
        create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH)
        print("\nAnalysis complete!")
        
    except Exception as e:
        print(f"\nError during analysis: {str(e)}")
        raise

if __name__ == "__main__":
    main()
