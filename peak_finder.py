import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from scipy.optimize import curve_fit
from scipy.signal import find_peaks
from scipy.fftpack import rfft, irfft, fftfreq, fft
from scipy.stats import linregress
import re
import json
from pptx import Presentation
from pptx.util import Inches

# Set a consistent style for all plots
plt.style.use('seaborn-v0_8-whitegrid')
plt.rcParams['figure.figsize'] = (12, 6)
plt.rcParams['font.size'] = 12
plt.rcParams['axes.titlesize'] = 14
plt.rcParams['axes.labelsize'] = 12
plt.rcParams['xtick.labelsize'] = 10
plt.rcParams['ytick.labelsize'] = 10

def extract_voltage_from_filename(filename):
    """Extract voltage value from filename."""
    match = re.search(r'Voltage_(-?\d+\.\d+)V_', filename)
    if match:
        return float(match.group(1))
    return None

def voltage_to_field(voltage, d=11.26):
    """Convert voltage to electric field in V/μm."""
    # d is in micrometers (μm)
    return voltage / d

def lorentzian(x, A, x0, gamma, B):
    """Calculate the Lorentzian function."""
    return ((A /np.pi) * (gamma / ((x - x0)**2 + gamma**2)) + B)

def most_common(lst):
    """Return most common element in a list."""
    return max(set(lst), key=lst.count)

def find_best_peak_range(data_directory, channel=1):
    """Find the best wavelength range based on peak analysis."""
    print("\nFinding optimal wavelength range for analysis...")
    
    # Find all CSV files
    csv_files = []
    for f in os.listdir(data_directory):
        if f.endswith('.csv'):
            voltage = extract_voltage_from_filename(f)
            if voltage is not None:
                csv_files.append((f, voltage))
    
    csv_files.sort(key=lambda x: x[1])
    print(f"Found {len(csv_files)} CSV files for peak analysis")
    
    # Track best peaks
    best_peaks = {
        'filename': [],
        'voltage': [],
        'wavelength': [],
        'extinction_ratio': [],
        'peak_index': [],
        'q_factor': []
    }
    
    # Scan all files
    for filename, voltage in csv_files:
        try:
            print(f"Processing {filename} (Voltage: {voltage}V)")
            filepath = os.path.join(data_directory, filename)
            df = pd.read_csv(filepath)
            
            # Ensure data is in the expected format
            if 'wavelength' not in df.columns or f'channel_{channel}' not in df.columns:
                print(f"  Warning: Expected columns not found in {filename}")
                continue
                
            x = np.array(df['wavelength'])
            y = np.array(df[f'channel_{channel}'])
            
            # Check if data needs to be converted to dBm
            if y.max() > 0:  # assuming linear scale
                y_db = 10 * np.log10(y)
            else:
                y_db = y
            
            # Find all dips (negative peaks)
            height = np.mean(y_db) - 1.5 * np.std(y_db)
            peaks, properties = find_peaks(-y_db, height=-height, distance=50, prominence=1.0)
            
            print(f"  Found {len(peaks)} potential peaks")
            
            # Process each peak
            for peak_idx in peaks:
                try:
                    # Extract region around peak
                    window = 100  # Points on each side
                    if peak_idx-window < 0 or peak_idx+window >= len(x):
                        continue
                    
                    xloc = x[peak_idx-window:peak_idx+window]
                    yloc = y_db[peak_idx-window:peak_idx+window]
                    
                    # Calculate extinction ratio directly
                    peak_min = np.min(yloc)
                    background = np.max(yloc)
                    ext_ratio = abs(peak_min - background)
                    
                    # Fit Lorentzian for better parameters
                    try:
                        p0 = [-(np.max(yloc)-np.min(yloc)),
                              xloc[round(len(xloc)/2)],
                              (xloc[-1]-xloc[0])/20,  # Initial guess for width
                              np.max(yloc)]
                        
                        params, _ = curve_fit(lorentzian, xloc, yloc, p0=p0,
                                           bounds=((-np.inf, xloc.min(), 0, -np.inf),
                                                  (0, xloc.max(), np.inf, np.inf)))
                        
                        # Calculate Q factor
                        Q = abs(params[1])/(abs(params[2])*2)
                        center_wavelength = params[1]
                    except Exception as fit_error:
                        print(f"    Lorentzian fit failed: {str(fit_error)}")
                        # If fitting fails, use direct measurements
                        Q = 0
                        center_wavelength = x[peak_idx]
                    
                    # Store this peak
                    best_peaks['filename'].append(filename)
                    best_peaks['voltage'].append(voltage)
                    best_peaks['wavelength'].append(center_wavelength)
                    best_peaks['extinction_ratio'].append(ext_ratio)
                    best_peaks['peak_index'].append(peak_idx)
                    best_peaks['q_factor'].append(Q)
                    
                    print(f"  Peak at {center_wavelength:.4f} nm, ER: {ext_ratio:.2f} dB, Q: {Q:.0f}")
                    
                except Exception as e:
                    print(f"  Error processing peak: {str(e)}")
                    continue
                    
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            continue
    
    # Create DataFrame for easier analysis
    peaks_df = pd.DataFrame(best_peaks)
    
    if len(peaks_df) == 0:
        print("No valid peaks found! Using default wavelength range.")
        return (1526.25, 1526.32), None  # Default range
    
    # Find the highest extinction ratio peak
    best_idx = peaks_df['extinction_ratio'].idxmax()
    best_peak = peaks_df.loc[best_idx]
    
    print("\nBest peak found:")
    print(f"File: {best_peak['filename']}")
    print(f"Voltage: {best_peak['voltage']}V")
    print(f"Wavelength: {best_peak['wavelength']:.4f} nm")
    print(f"Extinction Ratio: {best_peak['extinction_ratio']:.2f} dB")
    print(f"Q Factor: {best_peak['q_factor']:.0f}")
    
    # Create visualization of the best peak
    try:
        filename = best_peak['filename']
        filepath = os.path.join(data_directory, filename)
        df = pd.read_csv(filepath)
        
        x = np.array(df['wavelength'])
        y = np.array(df[f'channel_{channel}'])
        
        # Convert to dBm if necessary
        if y.max() > 0:
            y_db = 10 * np.log10(y)
        else:
            y_db = y
        
        # Find closest index to the peak wavelength
        peak_idx = np.abs(x - best_peak['wavelength']).argmin()
        
        # Extract region around peak for display
        window = 200  # Points on each side for wider view
        left_idx = max(0, peak_idx - window)
        right_idx = min(len(x) - 1, peak_idx + window)
        
        # Calculate suggested wavelength range (30% wider than visible peak)
        peak_min = np.min(y_db[left_idx:right_idx])
        background = np.max(y_db[left_idx:right_idx])
        half_depth = (background + peak_min) / 2
        
        # Find indices where the signal crosses half depth
        above_half = y_db[left_idx:right_idx] > half_depth
        transitions = np.where(np.diff(above_half))[0]
        
        if len(transitions) >= 2:
            left_edge = x[left_idx + transitions[0]]
            right_edge = x[left_idx + transitions[-1]]
            fwhm = right_edge - left_edge
            
            # Suggested range with 30% padding
            padding = fwhm * 0.3
            range_min = left_edge - padding
            range_max = right_edge + padding
        else:
            # Fallback if FWHM can't be determined
            wavelength_window = (x[right_idx] - x[left_idx])
            center = best_peak['wavelength']
            range_min = center - wavelength_window * 0.3
            range_max = center + wavelength_window * 0.3
        
        # Plot the peak
        plt.figure(figsize=(10, 6))
        plt.plot(x[left_idx:right_idx], y_db[left_idx:right_idx], 'b-')
        plt.axvline(best_peak['wavelength'], color='r', linestyle='--',
                   label=f'Peak: {best_peak["wavelength"]:.4f} nm')
        
        # Mark the suggested range
        plt.axvspan(range_min, range_max, alpha=0.2, color='green',
                   label=f'Analysis Range: {range_min:.4f} - {range_max:.4f} nm')
        
        plt.xlabel('Wavelength (nm)')
        plt.ylabel('Transmission (dBm)')
        plt.title(f'Best Peak Analysis (V={best_peak["voltage"]}V, ER={best_peak["extinction_ratio"]:.2f} dB)')
        plt.grid(True, alpha=0.3)
        plt.legend()
        
        plot_path = os.path.join(data_directory, 'best_peak_analysis.png')
        plt.savefig(plot_path, dpi=300)
        plt.close()
        
        print(f"Best peak analysis saved to: {plot_path}")
        print(f"Suggested wavelength range for analysis: {range_min:.4f} - {range_max:.4f} nm")
        
        return (range_min, range_max), best_peak
        
    except Exception as e:
        print(f"Error examining best peak: {str(e)}")
        # Use a reasonable default range
        peak_wavelength = best_peak['wavelength']
        return (peak_wavelength - 0.05, peak_wavelength + 0.05), best_peak

def plot_transmission_spectra(data_directory, wavelength_range, channel=1, save_dir=None):
    """Create overlay plot of transmission spectra at different voltages/fields."""
    print("\nCreating transmission spectra overlay plot...")
    
    fig, ax = plt.subplots(figsize=(12, 6))
    
    # Get all CSV files and sort by voltage
    csv_files = [(f, extract_voltage_from_filename(f))
                 for f in os.listdir(data_directory)
                 if f.endswith('.csv')]
    csv_files = [(f, v) for f, v in csv_files if v is not None]
    csv_files.sort(key=lambda x: x[1])
    
    # Create color map for different voltages
    colors = plt.cm.viridis(np.linspace(0, 1, len(csv_files)))
    
    for (filename, voltage), color in zip(csv_files, colors):
        # Load data
        df = pd.read_csv(os.path.join(data_directory, filename))
        wavelength = df['wavelength']
        transmission = df[f'channel_{channel}']
        
        # Convert to dBm if not already
        if transmission.max() > 0:  # assuming linear scale
            transmission = 10 * np.log10(transmission)
            
        # Calculate field
        field = voltage_to_field(voltage)
        
        ax.plot(wavelength, transmission, label=f'{field:.2f} V/μm', color=color, alpha=0.7)
    
    ax.set_xlabel('Wavelength (nm)')
    # Use the detected range or default if not provided
    xmin, xmax = wavelength_range if wavelength_range else (1526.25, 1526.32)
    ax.set_xlim(xmin, xmax)
    ax.set_ylabel('Transmission (dBm)')
    ax.set_title('Transmission Spectra at Different Electric Fields')
    ax.grid(True, alpha=0.3)
    
    # Create a legend outside the plot
    ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', title="Electric Field")
    plt.tight_layout()
    
    # Save plot
    plot_path = os.path.join(save_dir, 'transmission_spectra.png')
    plt.savefig(plot_path, bbox_inches='tight', dpi=300)
    plt.close()
    
    print(f"Transmission spectra plot saved to: {plot_path}")
    return plot_path

def analyze_voltage_series(data_directory, wavelength_range=None, channel=1):
    """Analyze voltage-dependent measurements from CSV files."""
    print(f"\nStarting analysis for channel_{channel}")
    print(f"Looking in directory: {data_directory}")
    
    voltage_data = {
        'voltages': [],
        'fields': [],
        'resonance_wavelengths': [],
        'q_factors': [],
        'dbm_values': [],
        'peak_indices': []
    }
    
    # Get all CSV files and sort by voltage
    csv_files = [(f, extract_voltage_from_filename(f))
                 for f in os.listdir(data_directory)
                 if f.endswith('.csv')]
    csv_files = [(f, v) for f, v in csv_files if v is not None]
    csv_files.sort(key=lambda x: x[1])
    
    print(f"\nFound {len(csv_files)} CSV files")
    for filename, voltage in csv_files:
        print(f"\nProcessing: {filename} (Voltage: {voltage}V)")
        
        try:
            # Load and process data
            path = os.path.join(data_directory, filename)
            df = pd.read_csv(path)
            print(f"DataFrame shape: {df.shape}")
            print(f"Columns: {df.columns.tolist()}")
            
            x = np.array(df['wavelength'])
            y = np.array(df[f'channel_{channel}'])
            print(f"Wavelength range: {x.min():.2f} to {x.max():.2f}")
            print(f"Signal range: {y.min():.2f} to {y.max():.2f}")
            
            # If wavelength range is provided, filter the data
            if wavelength_range:
                xmin, xmax = wavelength_range
                range_mask = (x >= xmin) & (x <= xmax)
                x_filtered = x[range_mask]
                y_filtered = y[range_mask]
                
                # If no data points in range, skip
                if len(x_filtered) == 0:
                    print(f"  No data points in wavelength range {xmin}-{xmax}, skipping")
                    continue
                    
                # Replace original data with filtered data
                x = x_filtered
                y = y_filtered
                print(f"  Filtered to {len(x)} points in range {xmin:.4f}-{xmax:.4f}")
            
            # Convert to dBm if not already
            if y.max() > 0:  # assuming linear scale
                y_db = 10 * np.log10(y)
            else:
                y_db = y
            
            # Initial peak finding
            height = np.mean(y_db) - 1.5 * np.std(y_db)  # Look for dips
            peaks, _ = find_peaks(-y_db, height=-height)
            print(f"Found {len(peaks)} initial peaks")
            
            if len(peaks) > 0:
                # Select the most prominent peak in the range
                peak_depths = [y_db[peak] for peak in peaks]
                main_peak_idx = peaks[np.argmin(peak_depths)]
                
                # Extract region around the peak
                window = min(100, len(x) // 4)  # Adaptive window size
                left_idx = max(0, main_peak_idx - window)
                right_idx = min(len(x) - 1, main_peak_idx + window)
                
                xloc = x[left_idx:right_idx]
                yloc = y_db[left_idx:right_idx]
                
                # Calculate extinction ratio directly
                peak_min = np.min(yloc)
                background = np.max(yloc)
                ext_ratio = abs(peak_min - background)
                
                try:
                    # Fit Lorentzian
                    p0 = [-(np.max(yloc)-np.min(yloc)),
                         xloc[len(xloc)//2],
                         (xloc[-1]-xloc[0])/10,
                         np.max(yloc)]
                    
                    params, covar = curve_fit(lorentzian, xloc, yloc, p0=p0,
                                            bounds=((-np.inf, xloc.min(), 0, -np.inf),
                                                   (0, xloc.max(), np.inf, np.inf)),
                                            maxfev=5000)
                    
                    # Calculate metrics
                    Q = abs(params[1])/(abs(params[2])*2)
                    resonance = params[1]
                    
                    # Calculate field from voltage
                    field = voltage_to_field(voltage)
                    
                    # Store results
                    voltage_data['voltages'].append(voltage)
                    voltage_data['fields'].append(field)
                    voltage_data['resonance_wavelengths'].append(resonance)
                    voltage_data['q_factors'].append(Q)
                    voltage_data['dbm_values'].append(ext_ratio)
                    voltage_data['peak_indices'].append(main_peak_idx)
                    print(f"  Peak analyzed: λ={resonance:.4f} nm, ER={ext_ratio:.2f} dB, Q={Q:.0f}")
                    
                except Exception as e:
                    print(f"  Error fitting peak: {str(e)}")
                    # Still store the data using direct measurements
                    resonance = x[main_peak_idx]
                    field = voltage_to_field(voltage)
                    
                    voltage_data['voltages'].append(voltage)
                    voltage_data['fields'].append(field)
                    voltage_data['resonance_wavelengths'].append(resonance)
                    voltage_data['q_factors'].append(0)  # Placeholder for failed fit
                    voltage_data['dbm_values'].append(ext_ratio)
                    voltage_data['peak_indices'].append(main_peak_idx)
                    print(f"  Peak recorded without fitting: λ={resonance:.4f} nm, ER={ext_ratio:.2f} dB")
            else:
                print("  No peaks found in this file")
                            
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            continue
    
    print("\nAnalysis complete:")
    print(f"Total peaks found: {len(voltage_data['voltages'])}")
    return voltage_data

def plot_voltage_dependencies(voltage_data, save_dir):
    """Create plots of parameters vs electric field."""
    print("\nCreating electric field dependency plots...")
    
    if len(voltage_data['voltages']) == 0:
        print("No data to plot!")
        return None
    
    fig = plt.figure(figsize=(15, 5))
    
    # Find reference wavelength (at or closest to 0V)
    zero_idx = None
    if 0 in voltage_data['voltages']:
        zero_idx = voltage_data['voltages'].index(0)
    else:
        # Find the closest to 0
        zero_idx = min(range(len(voltage_data['voltages'])),
                      key=lambda i: abs(voltage_data['voltages'][i]))
    
    ref_wavelength = voltage_data['resonance_wavelengths'][zero_idx] if zero_idx is not None else None
    
    # Extract relative wavelength shifts
    relative_wavelengths = []
    if ref_wavelength is not None:
        relative_wavelengths = [wl - ref_wavelength for wl in voltage_data['resonance_wavelengths']]
    else:
        relative_wavelengths = voltage_data['resonance_wavelengths']
    
    # Resonance Wavelength Shift vs Field
    ax1 = fig.add_subplot(131)
    scatter1 = ax1.scatter(voltage_data['fields'], relative_wavelengths,
                          alpha=0.2, c='blue', s=50, edgecolors='none')
    
    # Add trendline
    if len(voltage_data['fields']) > 1:
        # Linear regression
        slope, intercept, r_value, p_value, std_err = linregress(
            voltage_data['fields'], relative_wavelengths)
        
        x_line = np.array([min(voltage_data['fields']), max(voltage_data['fields'])])
        y_line = slope * x_line + intercept
        
        ax1.plot(x_line, y_line, 'r-', linewidth=2,
                label=f'Trend: {slope:.4f} nm/(V/μm)')
        
    ax1.set_xlabel('Electric Field (V/μm)')
    ax1.set_ylabel('Resonance Wavelength Shift (nm)')
    ax1.set_title('Resonance Wavelength Shift vs Electric Field')
    ax1.grid(True, alpha=0.3)
    ax1.axhline(y=0, color='k', linestyle='-', alpha=0.3)
    ax1.axvline(x=0, color='k', linestyle='-', alpha=0.3)
    ax1.legend()
    
    # Extinction Ratio vs Field
    ax2 = fig.add_subplot(132)
    scatter2 = ax2.scatter(voltage_data['fields'], voltage_data['dbm_values'],
                          alpha=0.2, c='green', s=50, edgecolors='none')
    
    # Add trendline
    if len(voltage_data['fields']) > 1:
        # Linear regression
        slope, intercept, r_value, p_value, std_err = linregress(
            voltage_data['fields'], voltage_data['dbm_values'])
        
        x_line = np.array([min(voltage_data['fields']), max(voltage_data['fields'])])
        y_line = slope * x_line + intercept
        
        ax2.plot(x_line, y_line, 'r-', linewidth=2,
                label=f'Trend: {slope:.4f} dB/(V/μm)')
    
    # Adjust y-axis range to focus on the data
    mean_ext = np.mean(voltage_data['dbm_values'])
    std_ext = np.std(voltage_data['dbm_values'])
    y_min = max(0, mean_ext - 3*std_ext)
    y_max = mean_ext + 3*std_ext
    ax2.set_ylim(y_min, y_max)
    
    ax2.set_xlabel('Electric Field (V/μm)')
    ax2.set_ylabel('Extinction Ratio (dB)')
    ax2.set_title('Extinction Ratio vs Electric Field')
    ax2.grid(True, alpha=0.3)
    ax2.axvline(x=0, color='k', linestyle='-', alpha=0.3)
    ax2.legend()
    
    # Q factor vs Field
    ax3 = fig.add_subplot(133)
    scatter3 = ax3.scatter(voltage_data['fields'], voltage_data['q_factors'],
                          alpha=0.2, c='purple', s=50, edgecolors='none')
    
    # Add trendline
    if len(voltage_data['fields']) > 1:
        # Filter out zero Q-factors (failed fits)
        valid_fields = [f for f, q in zip(voltage_data['fields'], voltage_data['q_factors']) if q > 0]
        valid_q = [q for q in voltage_data['q_factors'] if q > 0]
        
        if len(valid_fields) > 1:
            # Linear regression
            slope, intercept, r_value, p_value, std_err = linregress(
                valid_fields, valid_q)
            
            x_line = np.array([min(voltage_data['fields']), max(voltage_data['fields'])])
            y_line = slope * x_line + intercept
            
            ax3.plot(x_line, y_line, 'r-', linewidth=2,
                    label=f'Trend: {slope:.1f} /(V/μm)')
    
    ax3.set_xlabel('Electric Field (V/μm)')
    ax3.set_ylabel('Q Factor')
    ax3.set_title('Q Factor vs Electric Field')
    ax3.grid(True, alpha=0.3)
    ax3.axvline(x=0, color='k', linestyle='-', alpha=0.3)
    ax3.legend()
    
    plt.tight_layout()
    
    plot_path = os.path.join(save_dir, 'field_dependencies.png')
    plt.savefig(plot_path, dpi=300)
    plt.close()
    
    print(f"Plots saved to: {plot_path}")
    return plot_path

def create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH, SUMMARY_POINTS=''):
    """Create PowerPoint presentation with results."""
    print("\nCreating PowerPoint presentation...")
    
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = TITLE
    subtitle.text = AUTHOR
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = 'Electric Field-Dependent Analysis'
    
    # Create summary points from results
    if results and SUMMARY_POINTS == '':
        summary_text = "Analysis Results:\n\n"
        
        # Extract key insights
        for chip_id, data in results.items():
            summary_text += f"• Device: {chip_id}\n"
            
            if data['voltages']:
                # Wavelength shift summary
                start_wavelength = data['resonance_wavelengths'][0]
                end_wavelength = data['resonance_wavelengths'][-1]
                total_shift = end_wavelength - start_wavelength
                summary_text += f"• Total wavelength shift: {total_shift:.4f} nm\n"
                
                # Field range summary
                min_field = min(data['fields'])
                max_field = max(data['fields'])
                summary_text += f"• Field range: {min_field:.2f} to {max_field:.2f} V/μm\n"
                
                # Average extinction ratio
                avg_er = sum(data['extinction_ratios']) / len(data['extinction_ratios'])
                summary_text += f"• Average extinction ratio: {avg_er:.2f} dB\n"
                
                # Tuning efficiency (if multiple points)
                if len(data['fields']) > 1 and len(data['resonance_wavelengths']) > 1:
                    # Calculate slope using linear regression
                    slope, _, _, _, _ = linregress(data['fields'], data['resonance_wavelengths'])
                    summary_text += f"• Tuning efficiency: {slope:.4f} nm/(V/μm)\n"
                
                # Average Q factor (excluding zeros)
                valid_q = [q for q in data['q_factors'] if q > 0]
                if valid_q:
                    avg_q = sum(valid_q) / len(valid_q)
                    summary_text += f"• Average Q factor: {avg_q:.0f}\n"
    else:
        summary_text = SUMMARY_POINTS or "Analysis of resonance parameters vs applied electric field"
        
    content.text = summary_text
    
    # Add plots
    if image_paths:
        slide_layout = prs.slide_layouts[5]
        
        # First plot - field dependencies
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Field Dependencies"
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4)
        slide.shapes.add_picture(image_paths[0], left, top, width, height)
        
        # Second plot - transmission spectra
        if len(image_paths) > 1:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Transmission Spectra"
            slide.shapes.add_picture(image_paths[1], left, top, width, height)
        
        # Add best peak analysis plot if it exists
        best_peak_path = os.path.join(os.path.dirname(image_paths[0]), 'best_peak_analysis.png')
        if os.path.exists(best_peak_path):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Best Peak Analysis"
            slide.shapes.add_picture(best_peak_path, left, top, width, height)
    
    prs.save(PPTX_SAVE_PATH)
    print(f"Presentation saved to: {PPTX_SAVE_PATH}")

def analyze_data(DATA_DIRECTORY, TITLE, AUTHOR, wafer, reticle, chip, PD_Channel=1):
    """Main analysis function."""
    print(f"\nStarting analysis for {wafer}_{reticle}_{chip}")
    chip_ID = f"{wafer}_{reticle}_{chip}"
    
    # First, find the best wavelength range based on peak analysis
    wavelength_range, best_peak = find_best_peak_range(DATA_DIRECTORY, PD_Channel)
    
    # Analyze voltage series with the optimal wavelength range
    voltage_data = analyze_voltage_series(DATA_DIRECTORY, wavelength_range, PD_Channel)
    
    # Create voltage dependency plots
    field_plot_path = plot_voltage_dependencies(voltage_data, DATA_DIRECTORY)
    
    # Create transmission spectra overlay plot with optimal range
    spectra_plot_path = plot_transmission_spectra(DATA_DIRECTORY, wavelength_range, PD_Channel, DATA_DIRECTORY)
    
    # Create results dictionary
    results = {
        chip_ID: {
            'voltages': voltage_data['voltages'],
            'fields': voltage_data['fields'],
            'resonance_wavelengths': voltage_data['resonance_wavelengths'],
            'q_factors': voltage_data['q_factors'],
            'extinction_ratios': voltage_data['dbm_values']
        }
    }
    
    return results, [field_plot_path, spectra_plot_path]
    
def main():
    """Main execution function."""
    # === CONFIGURATION - EDIT THESE VALUES ===
    TITLE = 'Electric Field-Dependent Analysis'
    AUTHOR = 'aadarsh'
    wafer = "Sample"
    reticle = "Test"
    chip = "1"
    
    # Set this to the full path of your data folder
    DATA_DIRECTORY = '/Users/aadarsh/Desktop/memq/10V'
    
    # Where to save results
    PPTX_SAVE_PATH = 'field_analysis_results.pptx'
    JSON_SAVE_PATH = 'field_analysis_results.json'
    
    # Analysis settings
    PD_Channel = 1  # Which channel to analyze (1-4)
    # =====================================
    
    print("\nStarting electric field analysis script...")
    print(f"Data directory: {DATA_DIRECTORY}")
    print(f"Analyzing channel: {PD_Channel}")
    
    try:
        # Run analysis
        results, image_paths = analyze_data(DATA_DIRECTORY, TITLE, AUTHOR,
                                          wafer, reticle, chip, PD_Channel)
        
        # Save results to JSON
        with open(JSON_SAVE_PATH, 'w') as json_file:
            json.dump(results, json_file, indent=4)
        print(f"\nResults saved to: {JSON_SAVE_PATH}")
        
        # Create presentation
        create_presentation(TITLE, AUTHOR, results, image_paths, PPTX_SAVE_PATH)
        print("\nAnalysis complete!")
        
    except Exception as e:
        print(f"\nError during analysis: {str(e)}")
        import traceback
        traceback.print_exc()
        raise

if __name__ == "__main__":
    main()
